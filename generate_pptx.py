import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Colors
    BG_DARK = RGBColor(11, 13, 18)       # #0B0D12
    CARD_BG = RGBColor(20, 24, 34)       # #141822
    CARD_BORDER = RGBColor(212, 175, 55) # Gold
    GOLD = RGBColor(212, 175, 55)        # #D4AF37
    GOLD_LIGHT = RGBColor(245, 215, 110) # #F5D76E
    TEXT_WHITE = RGBColor(248, 250, 252) # #F8FAFC
    TEXT_MUTED = RGBColor(160, 174, 192) # #A0AEC0

    base_dir = os.path.dirname(os.path.abspath(__file__))
    hotel_dir = os.path.join(base_dir, "public", "assets", "hotel")
    output_path = os.path.join(base_dir, "Richman_Estate_Presentation.pptx")

    def set_slide_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.color.rgb = BG_DARK
        return bg

    def add_header(slide, title_text, category="RICHMAN ESTATE • DOSSIER OFFICIEL"):
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = GOLD

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = GOLD
        line.line.color.rgb = GOLD

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
        return card

    # ==========================================
    # SLIDE 1: COUVERTURE
    # ==========================================
    s1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s1)

    img_facade = os.path.join(hotel_dir, "01_facade_nuit.jpg")
    if not os.path.exists(img_facade):
        img_facade = os.path.join(hotel_dir, "01_facade_jour.jpg")
    if os.path.exists(img_facade):
        s1.shapes.add_picture(img_facade, Inches(7.0), Inches(1.2), Inches(5.5), Inches(5.5))

    add_card(s1, 0.8, 1.2, 5.8, 5.5, bg_color=CARD_BG, border_color=GOLD)

    badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.6), Inches(2.8), Inches(0.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = GOLD
    badge.line.color.rgb = GOLD
    tf_b = badge.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.text = "PROJET LÉGAL RP"
    p_b.font.size = Pt(11)
    p_b.font.bold = True
    p_b.font.color.rgb = BG_DARK
    p_b.alignment = PP_ALIGN.CENTER

    t_box = s1.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(5.0), Inches(1.8))
    tf = t_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "RICHMAN ESTATE"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    p2 = tf.add_paragraph()
    p2.text = "Hôtellerie de Prestige • Supercars • Événementiel"
    p2.font.size = Pt(13)
    p2.font.color.rgb = GOLD_LIGHT
    p2.space_before = Pt(8)

    d_box = s1.shapes.add_textbox(Inches(1.2), Inches(3.9), Inches(5.0), Inches(2.4))
    tf_d = d_box.text_frame
    tf_d.word_wrap = True
    
    items = [
        ("Localisation :", "Manoir de Richman (Richman Hills)"),
        ("Direction :", "Antonio Depresto (#3932) & Alex Breacker (#4526)"),
        ("Site Web :", "https://richman-estate.vercel.app/"),
        ("Discord :", "https://discord.gg/bdJhGdP3t9")
    ]
    for label, val in items:
        p = tf_d.add_paragraph() if tf_d.paragraphs[0].text else tf_d.paragraphs[0]
        p.text = f"{label} "
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = GOLD
        run = p.add_run()
        run.text = val
        run.font.size = Pt(11)
        run.font.bold = False
        run.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(3)

    # ==========================================
    # SLIDE 2: HISTOIRE ROLEPLAY
    # ==========================================
    s2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s2)
    add_header(s2, "Histoire Roleplay & Continuité du Domaine")

    col_w = 3.65
    top_pos = 1.8
    h_pos = 5.0

    cards_data = [
        ("1. LES ORIGINES (GENESIS)", 
         "Le domaine de Richman et le Golf Club sont des repères historiques pour Antonio Depresto et Alex Breacker.\n\n"
         "À l'époque, Antonio était employé chez Genesis Security, alors basée au Manoir de Richman. À sa fermeture, Alex Breacker a racheté l'ensemble de la flotte de sécurité et les droits pour y implanter Spartan Security."),
        
        ("2. L'ÈRE SPARTAN SECURITY", 
         "Antonio intègre Spartan Security auprès d'Alex. Grâce à sa rigueur et sa parfaite maîtrise du terrain, il devient rapidement cadre au sein de l'équipe d'une vingtaine d'agents.\n\n"
         "Le Manoir de Richman a été leur quartier général et le centre névralgique de leurs opérations pendant de longs mois."),
        
        ("3. RENAISSANCE : RICHMAN ESTATE", 
         "Après la restitution temporaire du domaine à l'État, Antonio et Alex s'associent d'égal à égal pour donner une nouvelle vie civile au domaine.\n\n"
         "Ils fondent Richman Estate : un complexe hôtelier prestigieux combiné à la mise en location de leur collection privée de supercars rarissimes.")
    ]

    for i, (title, content) in enumerate(cards_data):
        left_pos = 0.8 + i * (col_w + 0.38)
        add_card(s2, left_pos, top_pos, col_w, h_pos, CARD_BG, GOLD if i == 2 else CARD_BORDER)
        
        tb = s2.shapes.add_textbox(Inches(left_pos + 0.2), Inches(top_pos + 0.2), Inches(col_w - 0.4), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = GOLD

        tb_c = s2.shapes.add_textbox(Inches(left_pos + 0.2), Inches(top_pos + 1.0), Inches(col_w - 0.4), Inches(3.8))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True
        p_c = tf_c.paragraphs[0]
        p_c.text = content
        p_c.font.size = Pt(11)
        p_c.font.color.rgb = TEXT_WHITE

    # ==========================================
    # SLIDE 3: PÔLES D'ACTIVITÉ & SERVICES
    # ==========================================
    s3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s3)
    add_header(s3, "Description de l’Entreprise & Services Proposés")

    services = [
        ("HÔTELLERIE & RÉSIDENCES",
         "• Nuitées en chambres prestige & suites exécutives\n"
         "• Accès complet aux infrastructures de détente : piscine, courts de tennis, espaces lounges\n"
         "• Service de conciergerie et séjour d'affaires ou de villégiature pour les citoyens d'élite."),
        
        ("LOCATION DE SUPERCARS",
         "• Mise à disposition exclusive de supercars et sportives rares de collection\n"
         "• Véhicules dédiés aux cérémonies, mariages, tournages ou sorties d'exception\n"
         "• Encadrement 100% RP : contrat officiel, vérification du permis, caution obligatoire."),
        
        ("ÉVÉNEMENTIEL & RÉCEPTIONS",
         "• Privatisation totale ou partielle du domaine pour galas, mariages et séminaires d'entreprises\n"
         "• Organisation de réceptions publiques exclusives (soirées lounge, expositions automobiles)\n"
         "• Partenariats institutionnels avec le Gouvernement et la Justice.")
    ]

    for i, (title, content) in enumerate(services):
        left_pos = 0.8 + i * (col_w + 0.38)
        add_card(s3, left_pos, top_pos, col_w, h_pos, CARD_BG, GOLD)
        
        tb = s3.shapes.add_textbox(Inches(left_pos + 0.2), Inches(top_pos + 0.2), Inches(col_w - 0.4), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = GOLD

        tb_c = s3.shapes.add_textbox(Inches(left_pos + 0.2), Inches(top_pos + 1.1), Inches(col_w - 0.4), Inches(3.6))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True
        p_c = tf_c.paragraphs[0]
        p_c.text = content
        p_c.font.size = Pt(11)
        p_c.font.color.rgb = TEXT_WHITE

    # ==========================================
    # SLIDE 4: FLOTTE AUTOMOBILE & MODALITÉS
    # ==========================================
    s4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s4)
    add_header(s4, "Flotte Automobile d'Exception & Encadrement RP")

    add_card(s4, 0.8, 1.8, 5.6, 5.0, CARD_BG, GOLD)
    tb = s4.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = "CATALOGUE DE LA FLOTTE"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = GOLD

    tb_v = s4.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(5.2), Inches(3.9))
    tf_v = tb_v.text_frame
    tf_v.word_wrap = True
    tf_v.paragraphs[0].text = "🏎️ Supercars de Collection Privée :"
    tf_v.paragraphs[0].font.bold = True
    tf_v.paragraphs[0].font.size = Pt(12)
    tf_v.paragraphs[0].font.color.rgb = GOLD_LIGHT

    veh_list = [
        "• Benefactor Krieger — Hypercar allemande de prestige",
        "• Grotti Furia — Élégance et sonorité italienne pure",
        "• Principe Deveste Eight — Design futuriste d'exception",
        "",
        "🚘 Berlines & Véhicules d'Honneur :",
        "• Enus Super Diamond — Luxe britannique pour cortèges",
        "• Gallivanter Baller LWB — Confort SUV haut standing"
    ]
    for line in veh_list:
        p = tf_v.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        if "🚘" in line:
            p.font.bold = True
            p.font.color.rgb = GOLD_LIGHT
        else:
            p.font.color.rgb = TEXT_WHITE

    add_card(s4, 6.8, 1.8, 5.7, 5.0, CARD_BG, CARD_BORDER)
    tb_r = s4.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(0.6))
    p_r = tb_r.text_frame.paragraphs[0]
    p_r.text = "GESTION 100% RP DE LA LOCATION"
    p_r.font.size = Pt(15)
    p_r.font.bold = True
    p_r.font.color.rgb = GOLD

    tb_rc = s4.shapes.add_textbox(Inches(7.0), Inches(2.6), Inches(5.3), Inches(3.9))
    tf_rc = tb_rc.text_frame
    tf_rc.word_wrap = True
    rp_rules = [
        ("Propriété légitime :", "Tous les véhicules sont déjà acquis légalement par les gérants ou achetés via les concessions du serveur."),
        ("Prêt de clés :", "Effectué directement via les commandes de base du serveur (/donnercle, /givekey ou échange direct)."),
        ("Contrats & Cautions :", "Signature d'un contrat de location avec versement d'une caution obligatoire garantissant l'état du véhicule."),
        ("Restitution :", "État des lieux complet au retour et révocation des clés. En cas d'abus : dépôt de plainte RP et recours justice/police.")
    ]
    for k, v in rp_rules:
        p = tf_rc.add_paragraph() if tf_rc.paragraphs[0].text else tf_rc.paragraphs[0]
        p.text = f"✓ {k} "
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = GOLD
        run = p.add_run()
        run.text = v
        run.font.size = Pt(11)
        run.font.bold = False
        run.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(4)

    # ==========================================
    # SLIDE 5: TARIFS & PRESTATIONS
    # ==========================================
    s5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s5)
    add_header(s5, "Grille Tarifaire des Services & Prestations")

    tarifs = [
        ("HÉBERGEMENT", [
            ("Chambre Prestige", "1 500 $", "/ nuitée"),
            ("Suite Exécutive / Penthouse", "3 500 $", "/ nuitée"),
            ("Accès Détente & Piscine", "500 $", "/ personne / jour")
        ]),
        ("LOCATION VÉHICULES", [
            ("Location Supercar (Krieger/Furia)", "5 000 $ à 8 000 $", "/ jour"),
            ("Caution Obligatoire", "10 000 $ à 20 000 $", "(restituée au retour)"),
            ("Option Chauffeur Privé", "1 500 $", "/ événement")
        ]),
        ("PRIVATISATION & ÉVÉNEMENTS", [
            ("Privatisation Manoir (Soirée)", "15 000 $ à 25 000 $", "/ événement"),
            ("Réception / Gala d'Entreprise", "Sur Devis", "selon prestations"),
            ("Shooting & Tournage Photo/Vidéo", "3 000 $", "/ demi-journée")
        ])
    ]

    for i, (cat_title, rows) in enumerate(tarifs):
        left_pos = 0.8 + i * (col_w + 0.38)
        add_card(s5, left_pos, top_pos, col_w, h_pos, CARD_BG, GOLD)

        tb = s5.shapes.add_textbox(Inches(left_pos + 0.2), Inches(top_pos + 0.2), Inches(col_w - 0.4), Inches(0.6))
        p = tb.text_frame.paragraphs[0]
        p.text = cat_title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = GOLD

        tb_c = s5.shapes.add_textbox(Inches(left_pos + 0.2), Inches(top_pos + 0.9), Inches(col_w - 0.4), Inches(3.8))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True
        
        for name, price, sub in rows:
            p1 = tf_c.add_paragraph() if tf_c.paragraphs[0].text else tf_c.paragraphs[0]
            p1.text = name
            p1.font.size = Pt(11)
            p1.font.bold = True
            p1.font.color.rgb = TEXT_WHITE
            p1.space_before = Pt(6)

            p2 = tf_c.add_paragraph()
            p2.text = price
            p2.font.size = Pt(13)
            p2.font.bold = True
            p2.font.color.rgb = GOLD_LIGHT
            run = p2.add_run()
            run.text = f"  {sub}"
            run.font.size = Pt(10)
            run.font.bold = False
            run.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 6: ORGANIGRAMME & DRESS CODE
    # ==========================================
    s6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s6)
    add_header(s6, "Organigramme, Effectifs & Code Vestimentaire")

    add_card(s6, 0.8, 1.8, 5.6, 5.0, CARD_BG, GOLD)
    tb_o = s6.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(0.6))
    p_o = tb_o.text_frame.paragraphs[0]
    p_o.text = "STRUCTURE DES EFFECTIFS (< 30)"
    p_o.font.size = Pt(15)
    p_o.font.bold = True
    p_o.font.color.rgb = GOLD

    tb_oc = s6.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(5.2), Inches(4.0))
    tf_oc = tb_oc.text_frame
    tf_oc.word_wrap = True
    org_lines = [
        ("Direction Générale :", "PDG Antonio Depresto (#3932) & Co-PDG Alex Breacker (#4526)"),
        ("Pôle Hôtellerie :", "1x Responsable Accueil + 2x Réceptionnistes / Intendance"),
        ("Pôle Flotte :", "1x Responsable Flotte + 2x Préparateurs / Contrôleurs"),
        ("Pôle Événementiel :", "1x Chargé d'Événements & Relations Publiques"),
        ("Recrutement RP :", "Processus d'embauche sélectif avec contrat de travail RP.")
    ]
    for k, v in org_lines:
        p = tf_oc.add_paragraph() if tf_oc.paragraphs[0].text else tf_oc.paragraphs[0]
        p.text = f"• {k} "
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = GOLD_LIGHT
        run = p.add_run()
        run.text = v
        run.font.size = Pt(11)
        run.font.bold = False
        run.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(5)

    add_card(s6, 6.8, 1.8, 5.7, 5.0, CARD_BG, CARD_BORDER)
    tb_d = s6.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(0.6))
    p_d = tb_d.text_frame.paragraphs[0]
    p_d.text = "CODE VESTIMENTAIRE (DRESS CODE)"
    p_d.font.size = Pt(15)
    p_d.font.bold = True
    p_d.font.color.rgb = GOLD

    tb_dc = s6.shapes.add_textbox(Inches(7.0), Inches(2.6), Inches(5.3), Inches(4.0))
    tf_dc = tb_dc.text_frame
    tf_dc.word_wrap = True
    dress_lines = [
        ("👔 Direction :", "Costumes 3 pièces sur-mesure, montres de luxe, tenues de cérémonie."),
        ("🤵 Accueil & Réception :", "Costumes noirs cintrés, chemises blanches, cravates ou nœuds papillon soignés."),
        ("🚗 Pôle Flotte & Service :", "Chemises blanches élégantes, gilets de service noirs, pantalons habillés."),
        ("🌟 Cohérence visuelle :", "Tenues toujours impeccables reflétant le prestige et l'exclusivité du domaine.")
    ]
    for k, v in dress_lines:
        p = tf_dc.add_paragraph() if tf_dc.paragraphs[0].text else tf_dc.paragraphs[0]
        p.text = f"{k}\n"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = GOLD_LIGHT
        run = p.add_run()
        run.text = v
        run.font.size = Pt(11)
        run.font.bold = False
        run.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(5)

    # ==========================================
    # SLIDE 7: OBJECTIFS & SYNERGIES RP
    # ==========================================
    s7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s7)
    add_header(s7, "Objectifs Court/Long Terme & Synergies RP")

    add_card(s7, 0.8, 1.8, 5.6, 5.0, CARD_BG, GOLD)
    tb_obj = s7.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(0.6))
    p_obj = tb_obj.text_frame.paragraphs[0]
    p_obj.text = "OBJECTIFS DE DÉVELOPPEMENT"
    p_obj.font.size = Pt(15)
    p_obj.font.bold = True
    p_obj.font.color.rgb = GOLD

    tb_objc = s7.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(5.2), Inches(4.0))
    tf_objc = tb_objc.text_frame
    tf_objc.word_wrap = True
    objs = [
        ("Court terme :", "• Réalisation de la scène RP d'achat avec le Gouvernement\n• Mise en service du parc de supercars et contrats types\n• Recrutement et formation de l'équipe de réception\n• Grande soirée d'inauguration avec exposition automobile"),
        ("Long terme :", "• Devenir le pôle d'hôtellerie de luxe incontournable de Los Santos\n• Contrats réguliers de privatisation avec les institutions & entreprises\n• Concours d'élégance et rassemblements de supercars")
    ]
    for k, v in objs:
        p = tf_objc.add_paragraph() if tf_objc.paragraphs[0].text else tf_objc.paragraphs[0]
        p.text = f"{k}\n"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = GOLD_LIGHT
        run = p.add_run()
        run.text = v
        run.font.size = Pt(10)
        run.font.bold = False
        run.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(5)

    add_card(s7, 6.8, 1.8, 5.7, 5.0, CARD_BG, CARD_BORDER)
    tb_syn = s7.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(0.6))
    p_syn = tb_syn.text_frame.paragraphs[0]
    p_syn.text = "SYNERGIES & INTERACTIONS EN VILLE"
    p_syn.font.size = Pt(15)
    p_syn.font.bold = True
    p_syn.font.color.rgb = GOLD

    tb_sync = s7.shapes.add_textbox(Inches(7.0), Inches(2.6), Inches(5.3), Inches(4.0))
    tf_sync = tb_sync.text_frame
    tf_sync.word_wrap = True
    syn_list = [
        ("🏛️ Gouvernement & Justice :", "Accueil des séminaires d'État, galas de charité et réceptions officielles."),
        ("🔧 Garages & Concessions :", "Contrats réguliers d'entretien complet, révisions et pneumatiques pour la flotte."),
        ("📰 Weazel News & Médias :", "Campagnes publicitaires, couverture des soirées mondaines et interviews."),
        ("🍸 Commerces & Traiteurs :", "Commandes locales régulières pour les buffets, bars et réceptions.")
    ]
    for k, v in syn_list:
        p = tf_sync.add_paragraph() if tf_sync.paragraphs[0].text else tf_sync.paragraphs[0]
        p.text = f"{k} "
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = GOLD_LIGHT
        run = p.add_run()
        run.text = v
        run.font.size = Pt(10)
        run.font.bold = False
        run.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(5)

    # ==========================================
    # SLIDE 8: ENGAGEMENTS HRP & NOTE STAFF
    # ==========================================
    s8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s8)
    add_header(s8, "Engagements HRP & Conformité au Règlement")

    add_card(s8, 0.8, 1.8, 11.733, 5.0, CARD_BG, GOLD)
    
    tb_hrp = s8.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(11.0), Inches(4.3))
    tf_hrp = tb_hrp.text_frame
    tf_hrp.word_wrap = True

    engagements = [
        ("0 SCRIPT / 0 DÉVELOPPEMENT DEMANDÉ :", "L'entreprise est 100% autonome et utilise uniquement les mécaniques natives du serveur (commandes de clés existantes, factures, virements bancaires, contrats RP)."),
        ("0 MAPPING / 0 WHITELIST SPÉCIALE :", "Le domaine sera exploité dans son état naturel sur la carte, sans aucune charge technique pour les équipes du serveur."),
        ("0 DEMANDE D'ARMES OU DE VÉHICULES DE SPAWN :", "L'ensemble de la flotte de supercars et de berlines est issu des biens propres légitimes des gérants déjà enregistrés en ville."),
        ("ACQUISITION DU LIEU EN SCÈNE RP :", "L'accès au domaine sera négocié et acheté lors d'un rendez-vous RP officiel avec le Gouvernement de Los Santos."),
        ("RESPECT STRICT DES QUOTAS :", "Effectifs maintenus à taille humaine (< 30 membres), axés sur la qualité du Roleplay, l'accueil et l'animation de la ville.")
    ]

    for title, desc in engagements:
        p = tf_hrp.add_paragraph() if tf_hrp.paragraphs[0].text else tf_hrp.paragraphs[0]
        p.text = f"✅ {title} "
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = GOLD
        run = p.add_run()
        run.text = desc
        run.font.size = Pt(11)
        run.font.bold = False
        run.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(8)

    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_presentation()
